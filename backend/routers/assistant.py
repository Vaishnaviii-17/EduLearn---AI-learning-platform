from fastapi import APIRouter, UploadFile, File, HTTPException, Body
import os
import re
import json
from dotenv import load_dotenv
import google.generativeai as genai
from transformers import pipeline
import PyPDF2
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

from fastapi.middleware.cors import CORSMiddleware

# Allow React frontend
origins = [
    "http://localhost:3000",
    "http://127.0.0.1:3000"
    "https://edulearn-ai-learning-platform-1.onrender.com"
]

def setup_cors(app):
    app.add_middleware(
        CORSMiddleware,
        allow_origins=origins,
        allow_credentials=True,
        allow_methods=["*"],
        allow_headers=["*"],
    )


# ----------------------------
# Load environment variables
# ----------------------------
load_dotenv()
router = APIRouter(prefix="/assistant", tags=["assistant"])

# ----------------------------
# Configure Gemini API
# ----------------------------
GEMINI_API_KEY = os.getenv("GOOGLE_API_KEY")
if not GEMINI_API_KEY:
    raise RuntimeError("❌ GOOGLE_API_KEY not found in .env file.")
genai.configure(api_key=GEMINI_API_KEY)

# ----------------------------
# Fallback models (HuggingFace)
# ----------------------------
quick_summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
detailed_summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

# ----------------------------
# File extractors
# ----------------------------
def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return " ".join([page.extract_text() or "" for page in reader.pages])

def extract_text_from_docx(file):
    document = docx.Document(file)
    return " ".join([p.text for p in document.paragraphs])

def extract_text_from_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return " ".join(text)

def extract_text_from_image(file):
    return pytesseract.image_to_string(Image.open(file))

def extract_text(upload_file: UploadFile):
    fname = upload_file.filename.lower()
    fobj = upload_file.file
    if fname.endswith(".pdf"):
        return extract_text_from_pdf(fobj)
    if fname.endswith((".docx", ".doc")):
        return extract_text_from_docx(fobj)
    if fname.endswith((".pptx", ".ppt")):
        return extract_text_from_pptx(fobj)
    if fname.endswith((".png", ".jpg", ".jpeg", ".tif", ".tiff")):
        return extract_text_from_image(fobj)
    raise HTTPException(status_code=400, detail="Unsupported file type")

# ----------------------------
# Helpers
# ----------------------------
def clean_text(text: str) -> str:
    text = re.sub(r"-\s+", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()

# ----------------------------
# Summarization (Gemini + HF fallback)
# ----------------------------
def summarize_with_gemini(text: str, mode="quick"):
    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        if mode == "quick":
            prompt = (
                "Analyze the document and list ONLY the main topics or key sections "
                "as concise bullet points with one-line student-friendly explanations.\n\n"
            )
        else:
            prompt = (
                "Provide a detailed, topic-wise student-friendly summary of this document. "
                "Use headings and bullet points, be clear and structured.\n\n"
            )
        response = model.generate_content(prompt + text[:12000])
        if response and getattr(response, "text", None):
            return response.text.strip()
        return None
    except Exception as e:
        print("⚠️ Gemini summarization failed:", e)
        return None

def summarize_with_huggingface(text: str, mode="quick"):
    text = clean_text(text)
    max_chunk = 600 if mode == "quick" else 800
    words = text.split()
    chunks, cur = [], []
    for w in words:
        cur.append(w)
        if len(cur) >= max_chunk:
            chunks.append(" ".join(cur))
            cur = []
    if cur:
        chunks.append(" ".join(cur))
    summarizer = quick_summarizer if mode == "quick" else detailed_summarizer
    summaries = []
    for c in chunks[:3]:
        try:
            out = summarizer(
                c,
                max_length=150 if mode == "quick" else 200,
                min_length=30 if mode == "quick" else 60,
                do_sample=False,
            )
            summaries.append(out[0]["summary_text"])
        except Exception as e:
            print("⚠️ HF summarizer failed:", e)
            summaries.append(c[:300])
    return "\n\n".join(summaries)

# ----------------------------
# Flowchart generator (replaces Flashcards)
# ----------------------------
def generate_flowchart(text: str):
    """
    Generate a JSON-structured flowchart with nodes and edges.
    Handles non-standard Gemini responses gracefully.
    """
    import json
    import re
    from google.api_core.exceptions import GoogleAPIError

    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")

        prompt = (
            "Create a step-by-step conceptual flowchart based on the document below. "
            "Output ONLY a valid JSON with this structure:\n"
            "{"
            "  'nodes': [{'id': '1', 'label': 'Start'}, ...], "
            "  'edges': [{'source': '1', 'target': '2'}, ...]"
            "}\n"
            "Do not include any commentary, markdown, or text outside the JSON."
        )

        response = model.generate_content(prompt + "\n\n" + text[:12000])

        if not response or not getattr(response, "text", None):
            return {"nodes": [], "edges": []}

        txt = response.text.strip()

        # 🧹 Fix common formatting issues before parsing
        txt = txt.replace("'", '"')  # convert single quotes to double quotes
        txt = re.sub(r"```(?:json)?|```", "", txt)  # remove code fences
        txt = txt.strip()

        # Extract JSON substring if Gemini added commentary
        json_part = None
        try:
            json_part = txt[txt.index("{"): txt.rindex("}") + 1]
        except Exception:
            json_part = txt

        # ✅ Try to parse safely
        data = json.loads(json_part)

        if isinstance(data, dict) and "nodes" in data:
            return data

        # Fallback empty
        return {"nodes": [], "edges": []}

    except (json.JSONDecodeError, GoogleAPIError, Exception) as e:
        print("⚠️ Flowchart generation failed:", e)
        return {"nodes": [], "edges": []}

# ----------------------------
# Quiz generator + graders
# ----------------------------
def generate_quiz(text: str):
    try:
        model = genai.GenerativeModel("models/gemini-2.5-flash")
        prompt = (
            "Create a 10-question multiple-choice quiz from the text below. "
            "Output as strict JSON array: "
            "[{\"question\":\"...\",\"options\":{\"A\":\"..\",\"B\":\"..\",\"C\":\"..\",\"D\":\"..\"},\"answer\":\"B\"},...]\n\n"
        )
        response = model.generate_content(prompt + text[:12000])
        if response and getattr(response, "text", None):
            txt = response.text.strip()
            json_text = txt[txt.index("[") : txt.rindex("]") + 1]
            parsed = json.loads(json_text)
            quiz = []
            for item in parsed[:10]:
                q = item.get("question", "").strip()
                opts = item.get("options", {})
                ans = item.get("answer", "").strip().upper()
                normalized = {k: opts.get(k, "").strip() for k in ["A", "B", "C", "D"]}
                if q:
                    quiz.append({"question": q, "options": normalized, "answer": ans})
            return quiz
    except Exception as e:
        print("⚠️ Quiz generation failed:", e)
    return []

def grade_quiz(submitted):
    results, score, total = [], 0, 0
    try:
        if isinstance(submitted, dict) and "quiz" in submitted and "answers" in submitted:
            quiz = submitted["quiz"]
            answers = submitted["answers"]
            for i, item in enumerate(quiz):
                total += 1
                correct = item.get("answer", "").upper().strip()
                sel = answers.get(str(i), "").upper()
                ok = sel == correct
                if ok:
                    score += 1
                results.append({"index": i, "selected": sel, "correct": correct, "ok": ok})
    except Exception as e:
        print("⚠️ Error grading quiz:", e)
        raise
    return {"score": score, "total": total, "results": results}

# ----------------------------
# Endpoints
# ----------------------------
@router.post("/summarize/quick")
async def summarize_quick(file: UploadFile = File(...)):
    text = extract_text(file)
    if not text.strip():
        raise HTTPException(status_code=400, detail="No text extracted")
    summary = summarize_with_gemini(text, "quick") or summarize_with_huggingface(text, "quick")
    return {"summary": summary or "No summary available."}

@router.post("/summarize/detailed")
async def summarize_detailed(file: UploadFile = File(...)):
    text = extract_text(file)
    if not text.strip():
        raise HTTPException(status_code=400, detail="No text extracted")
    summary = summarize_with_gemini(text, "detailed") or summarize_with_huggingface(text, "detailed")
    return {"summary": summary or "No summary available."}

@router.post("/flowchart")
async def create_flowchart(file: UploadFile = File(...)):
    text = extract_text(file)
    if not text.strip():
        raise HTTPException(status_code=400, detail="No text extracted from file")
    flowchart = generate_flowchart(text)
    return {"flowchart": flowchart}


@router.post("/quiz")
async def create_quiz(file: UploadFile = File(...)):
    text = extract_text(file)
    if not text.strip():
        raise HTTPException(status_code=400, detail="No text extracted from file")
    quiz = generate_quiz(text)
    return {"quiz": quiz}

@router.post("/quiz/submit")
async def submit_quiz(payload: dict = Body(...)):
    try:
        result = grade_quiz(payload)
        return result
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))
